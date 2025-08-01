import os
import uuid
from typing import List, Tuple
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import base64
from io import BytesIO

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = {'.pdf', '.docx', '.pptx'}
    
    def process_document(self, file_path: str, filename: str) -> Tuple[List[str], str]:
        """Process a document and return chunks of text and file type"""
        file_extension = Path(filename).suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        if file_extension == '.pdf':
            return self._process_pdf(file_path), 'pdf'
        elif file_extension == '.docx':
            return self._process_docx(file_path), 'docx'
        elif file_extension == '.pptx':
            return self._process_pptx(file_path), 'pptx'
    
    def _process_pdf(self, file_path: str) -> List[str]:
        """Extract text from PDF file"""
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        # Split into smaller chunks if page is too long
                        page_chunks = self._split_text(text, max_length=1000)
                        for i, chunk in enumerate(page_chunks):
                            chunks.append(f"Page {page_num + 1}, Chunk {i + 1}: {chunk}")
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")
        
        return chunks
    
    def _process_docx(self, file_path: str) -> List[str]:
        """Extract text from DOCX file"""
        chunks = []
        try:
            doc = Document(file_path)
            current_chunk = ""
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    current_chunk += text + "\n"
                    
                    # Create chunk when it reaches reasonable size
                    if len(current_chunk) > 800:
                        chunks.append(current_chunk.strip())
                        current_chunk = ""
            
            # Add remaining text as final chunk
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
                
        except Exception as e:
            raise Exception(f"Error processing DOCX: {str(e)}")
        
        return chunks
    
    def _process_pptx(self, file_path: str) -> List[str]:
        """Extract text from PPTX file"""
        chunks = []
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip() != f"Slide {slide_num + 1}:":
                    chunks.append(slide_text.strip())
                    
        except Exception as e:
            raise Exception(f"Error processing PPTX: {str(e)}")
        
        return chunks
    
    def _split_text(self, text: str, max_length: int = 1000) -> List[str]:
        """Split text into smaller chunks"""
        if len(text) <= max_length:
            return [text]
        
        chunks = []
        sentences = text.split('. ')
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk + sentence) <= max_length:
                current_chunk += sentence + ". "
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + ". "
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def process_image(self, image_data: str) -> str:
        """Process base64 image data and return description"""
        try:
            # Decode base64 image
            image_bytes = base64.b64decode(image_data.split(',')[1] if ',' in image_data else image_data)
            image = Image.open(BytesIO(image_bytes))
            
            # For now, return basic image info
            # In a full implementation, you'd use a vision model here
            return f"Image uploaded: {image.format} format, size: {image.size}"
            
        except Exception as e:
            return f"Error processing image: {str(e)}"
